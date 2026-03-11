from typing import List

import httpx
from bs4 import BeautifulSoup

from ..config import get_settings
from ..models import Chunk, Document, DocumentContent
from ..vectorstore import upsert_chunks
from . import embedding_service, web_service

settings = get_settings()


def _chunk_text(text: str, max_chars: int = 800) -> List[str]:
    chunks = []
    current = []
    current_len = 0
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        if current_len + len(line) > max_chars and current:
            chunks.append("\n".join(current))
            current = [line]
            current_len = len(line)
        else:
            current.append(line)
            current_len += len(line)
    if current:
        chunks.append("\n".join(current))
    return chunks


# ---------------------------------------------------------------------------
# 多格式文件文本提取
# ---------------------------------------------------------------------------

def _extract_text_from_file(path: str, doc_type: str) -> str:
    """根据文档类型使用对应库提取文本。"""
    doc_type = doc_type.lower()

    if doc_type in ("doc", "docx"):
        return _extract_docx(path)
    elif doc_type == "pdf":
        return _extract_pdf(path)
    elif doc_type in ("xls", "xlsx"):
        return _extract_excel(path)
    elif doc_type in ("ppt", "pptx"):
        return _extract_pptx(path)
    elif doc_type == "txt":
        return _extract_plain(path)
    else:
        return _extract_plain(path)


def _extract_docx(path: str) -> str:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except Exception:
        return _extract_plain(path)


def _extract_pdf(path: str) -> str:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
                # 提取表格
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        cells = [str(c).strip() for c in row if c]
                        if cells:
                            texts.append(" | ".join(cells))
        return "\n".join(texts)
    except Exception:
        return _extract_plain(path)


def _extract_excel(path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            texts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None]
                if cells:
                    texts.append(" | ".join(cells))
        wb.close()
        return "\n".join(texts)
    except Exception:
        return _extract_plain(path)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            texts.append(f"[Slide {slide_idx}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
        return "\n".join(texts)
    except Exception:
        return _extract_plain(path)


def _extract_plain(path: str) -> str:
    try:
        with open(path, "rb") as f:
            raw = f.read()
        return raw.decode("utf-8", errors="ignore")
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# 核心解析流程
# ---------------------------------------------------------------------------

def chunk_and_index_text(db, document: Document, text: str) -> None:
    chunks_text = _chunk_text(text)
    
    # 获取有效的不为空的 chunks
    valid_chunks_text = [c for c in chunks_text if c and c.strip()]
    
    if not valid_chunks_text:
        document.status = "failed"
        db.add(document)
        db.commit()
        return

    embeddings = embedding_service.embed_texts(valid_chunks_text)

    db_chunks: List[Chunk] = []
    for idx, c in enumerate(valid_chunks_text):
        db_chunk = Chunk(
            document_id=document.id,
            content=c,
            position=idx,
            section_path=None,
            modality="text",
            layer="detail",
            embedding_index="global",
        )
        db.add(db_chunk)
        db_chunks.append(db_chunk)
    db.commit()

    ids = [f"chunk-{document.id}-{c.id}" for c in db_chunks]
    metadatas = [
        {
            "document_id": c.document_id,
            "chunk_id": c.id,
            "user_id": document.user_id,
            "modality": c.modality,
            "layer": c.layer,
            "scope": "global",
        }
        for c in db_chunks
    ]
    upsert_chunks(ids=ids, embeddings=embeddings, metadatas=metadatas, documents=chunks_text)

    document.status = "parsed"
    db.add(document)
    db.commit()


def parse_text_document(db, document: Document) -> None:
    if not document.original_path:
        return
    path = document.original_path

    text = _extract_text_from_file(path, document.doc_type)
    if not text.strip():
        document.status = "failed"
        db.add(document)
        db.commit()
        return

    content = DocumentContent(
        document_id=document.id,
        raw_text=text,
        cleaned_text=text,
    )
    db.add(content)
    db.commit()

    chunk_and_index_text(db, document, text)


def fetch_and_parse_web(db, document: Document, url: str) -> None:
    try:
        # 优先使用阿里云网页解析 API 获取高质量正文
        text = web_service.scrape_page(url)
        raw_html = f"[Aliyun Scrape] {url}"
        
        # 如果解析失败或者内容为空，降级到默认抓取
        if not text.strip():
            logger.warning(f"阿里云网页解析为空，降级到普通抓取: {url}")
            resp = httpx.get(url, timeout=15, follow_redirects=True)
            resp.raise_for_status()
            raw_html = resp.text
            
            soup = BeautifulSoup(raw_html, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            
            if not text.strip():
                text = raw_html

    except Exception as e:
        logger.error(f"网页抓取失败: {e}")
        document.status = "failed"
        db.add(document)
        db.commit()
        return

    content = DocumentContent(
        document_id=document.id,
        raw_text=raw_html,
        cleaned_text=text,
    )
    db.add(content)
    
    # 更新文档状态为已解析
    document.status = "parsed"
    db.add(document)
    
    db.commit()

    chunk_and_index_text(db, document, text)
